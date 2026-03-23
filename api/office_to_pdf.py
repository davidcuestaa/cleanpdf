"""
API endpoint: POST /api/office_to_pdf
Convierte DOCX, XLSX, PPTX a PDF.

Entorno: Vercel serverless (Python puro, sin binarios del sistema).
Librerías: python-docx, openpyxl, python-pptx, reportlab

Fixes v3:
- Fuente DejaVuSans (incluida en reportlab) para soporte Unicode completo:
  bullets •, guiones –, flechas →, comillas tipográficas "", etc.
- Fallback seguro a Helvetica con sustitución de caracteres si DejaVu no carga
- Tamaños de fuente, colores, alineaciones y márgenes leídos del documento real
- Tablas en orden correcto (intercaladas con párrafos, no al final)
- Cabeceras de tabla con color real del documento si está disponible
"""
import os
import io
import cgi
import json
from http.server import BaseHTTPRequestHandler

try:
    from docx import Document
    from docx.oxml.ns import qn as wqn
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    import openpyxl
    XLSX_OK = True
except ImportError:
    XLSX_OK = False

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from reportlab.lib.pagesizes import A4, landscape as rl_landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Table, TableStyle, PageBreak, HRFlowable)
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    RL_OK = True
except ImportError:
    RL_OK = False

SUPPORTED = {".docx", ".xlsx", ".pptx"}


# ── Registro de fuente Unicode ────────────────────────────────────────────────
# DejaVuSans viene incluida con reportlab — soporta •, –, →, ©, tildes, etc.

_FONTS_READY = False

def _ensure_fonts():
    global _FONTS_READY
    if _FONTS_READY or not RL_OK:
        return
    try:
        import reportlab
        rl_dir = os.path.dirname(reportlab.__file__)
        candidates = [
            (os.path.join(rl_dir, "fonts", "DejaVuSans.ttf"),
             os.path.join(rl_dir, "fonts", "DejaVuSans-Bold.ttf")),
            (os.path.join(rl_dir, "fonts", "dejavu", "DejaVuSans.ttf"),
             os.path.join(rl_dir, "fonts", "dejavu", "DejaVuSans-Bold.ttf")),
            ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
             "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        ]
        for reg_path, bold_path in candidates:
            if os.path.isfile(reg_path):
                pdfmetrics.registerFont(TTFont("DVS", reg_path))
                if os.path.isfile(bold_path):
                    pdfmetrics.registerFont(TTFont("DVS-Bold", bold_path))
                else:
                    pdfmetrics.registerFont(TTFont("DVS-Bold", reg_path))
                _FONTS_READY = True
                break
    except Exception:
        pass
    _FONTS_READY = True  # no reintentar aunque falle


def _fn(bold=False, italic=False):
    """Nombre de fuente: DejaVuSans si está disponible, Helvetica si no."""
    if _FONTS_READY and "DVS" in (pdfmetrics.getRegisteredFontNames()
                                   if hasattr(pdfmetrics, "getRegisteredFontNames")
                                   else []):
        return "DVS-Bold" if bold else "DVS"
    # Fallback Helvetica
    if bold and italic: return "Helvetica-BoldOblique"
    if bold:            return "Helvetica-Bold"
    if italic:          return "Helvetica-Oblique"
    return "Helvetica"


def _safe(text):
    """
    Escapa XML para ReportLab.
    Si no hay fuente Unicode, sustituye caracteres problemáticos.
    """
    t = (str(text)
         .replace("&", "&amp;").replace("<", "&lt;")
         .replace(">", "&gt;").replace('"', "&quot;"))

    has_unicode_font = _FONTS_READY and "DVS" in (
        pdfmetrics.getRegisteredFontNames()
        if hasattr(pdfmetrics, "getRegisteredFontNames") else [])

    if not has_unicode_font:
        subs = {
            "\u2022": "-", "\u2023": "-", "\u2043": "-",
            "\u25cf": "-", "\u25e6": "-", "\u00b7": "-",
            "\u2013": "-", "\u2014": "--", "\u2012": "-",
            "\u2018": "'", "\u2019": "'",
            "\u201c": '"', "\u201d": '"',
            "\u2026": "...", "\u00a0": " ",
            "\u2192": "->", "\u2190": "<-",
            "\u00ae": "(R)", "\u00a9": "(C)", "\u2122": "(TM)",
        }
        for orig, repl in subs.items():
            t = t.replace(orig, repl)
    return t


def _color(rgb_str):
    h = str(rgb_str).lstrip("#")
    if len(h) == 6:
        try:
            return colors.HexColor(f"#{h}")
        except Exception:
            pass
    return None


def _parse_multipart(handler):
    ct = handler.headers.get("Content-Type", "")
    if "multipart/form-data" not in ct:
        return None, None, "Content-Type must be multipart/form-data"
    env = {"REQUEST_METHOD": "POST", "CONTENT_TYPE": ct,
           "CONTENT_LENGTH": handler.headers.get("Content-Length", "0")}
    body = handler.rfile.read(int(handler.headers.get("Content-Length", 0)))
    fs = cgi.FieldStorage(fp=io.BytesIO(body), environ=env, keep_blank_values=True)
    if "file" not in fs:
        return None, None, "No 'file' field in request"
    f = fs["file"]
    return f.file.read(), f.filename, None


# ─────────────────────────────────────────────────────────────────────────────
# DOCX → PDF
# ─────────────────────────────────────────────────────────────────────────────

def docx_to_pdf(file_bytes):
    if not (DOCX_OK and RL_OK):
        raise RuntimeError("python-docx y reportlab son necesarios")

    _ensure_fonts()
    doc = Document(io.BytesIO(file_bytes))
    buf = io.BytesIO()

    # Tamaño de página y márgenes reales del documento
    try:
        sec = doc.sections[0]
        def emu2cm(e): return e / 914400 * 2.54 * cm
        pw = max(min(emu2cm(sec.page_width),  30*cm), 10*cm)
        ph = max(min(emu2cm(sec.page_height), 42*cm), 10*cm)
        lm = max(min(emu2cm(sec.left_margin),  5*cm), 0.5*cm)
        rm = max(min(emu2cm(sec.right_margin), 5*cm), 0.5*cm)
        tm = max(min(emu2cm(sec.top_margin),   5*cm), 0.5*cm)
        bm = max(min(emu2cm(sec.bottom_margin),5*cm), 0.5*cm)
    except Exception:
        pw, ph = A4
        lm = rm = 2.5*cm
        tm = bm = 2.5*cm

    pdf   = SimpleDocTemplate(buf, pagesize=(pw, ph),
                              leftMargin=lm, rightMargin=rm,
                              topMargin=tm, bottomMargin=bm)
    base  = getSampleStyleSheet()
    story = []

    H_SIZES   = {1: 22, 2: 18, 3: 15, 4: 13, 5: 12, 6: 11}
    ALIGN_MAP = {"LEFT": TA_LEFT, "CENTER": TA_CENTER,
                 "RIGHT": TA_RIGHT, "JUSTIFY": TA_JUSTIFY,
                 "BOTH": TA_JUSTIFY, None: TA_LEFT}

    def get_align(para):
        try:
            return ALIGN_MAP.get(para.alignment.name if para.alignment else None, TA_LEFT)
        except Exception:
            return TA_LEFT

    def get_indent(pf):
        try:
            return float(pf.left_indent.pt) if pf.left_indent else 0
        except Exception:
            return 0

    def get_spacing(pf):
        try:
            sb = float(pf.space_before.pt) if pf.space_before else 0
        except Exception:
            sb = 0
        try:
            sa = float(pf.space_after.pt) if pf.space_after else 4
        except Exception:
            sa = 4
        return sb, sa

    def make_style(para):
        sname  = para.style.name if para.style else "Normal"
        pf     = para.paragraph_format
        align  = get_align(para)
        sb, sa = get_spacing(pf)
        indent = get_indent(pf)

        if sname.startswith("Heading"):
            try:
                lvl = int(sname.split()[-1])
            except Exception:
                lvl = 1
            sz = H_SIZES.get(lvl, 12)
            return ParagraphStyle(
                f"h{lvl}_{id(para)}", parent=base["Normal"],
                fontSize=sz, fontName=_fn(bold=True),
                spaceBefore=max(sb, sz * 0.6), spaceAfter=max(sa, 4),
                leading=sz * 1.3, alignment=align, leftIndent=indent,
                textColor=colors.HexColor("#111827"),
            )

        font_size = 11.0
        txt_color = colors.HexColor("#1a1a1a")
        bold = italic = False

        if para.runs:
            sizes = [r.font.size.pt for r in para.runs if r.font.size]
            if sizes:
                font_size = max(6.0, min(72.0, sorted(sizes)[len(sizes)//2]))
            bold   = sum(1 for r in para.runs if r.bold)   > len(para.runs) * 0.5
            italic = sum(1 for r in para.runs if r.italic) > len(para.runs) * 0.5
            for r in para.runs:
                try:
                    if r.font.color and r.font.color.type is not None:
                        c = _color(str(r.font.color.rgb))
                        if c:
                            txt_color = c
                            break
                except Exception:
                    pass

        return ParagraphStyle(
            f"p_{id(para)}", parent=base["Normal"],
            fontName=_fn(bold=bold, italic=italic),
            fontSize=font_size, leading=font_size * 1.35,
            spaceAfter=max(sa, 2), spaceBefore=sb,
            alignment=align, textColor=txt_color, leftIndent=indent,
        )

    def rich_text(para):
        if not para.runs:
            return _safe(para.text)
        parts = []
        for run in para.runs:
            t = _safe(run.text)
            if not t:
                continue
            try:
                if run.font.size:
                    t = f'<font size="{run.font.size.pt:.1f}">{t}</font>'
            except Exception:
                pass
            try:
                if run.font.color and run.font.color.type is not None:
                    c = _color(str(run.font.color.rgb))
                    if c:
                        t = f'<font color="{c.hexval()}">{t}</font>'
            except Exception:
                pass
            b, i, u = run.bold, run.italic, run.underline
            if b and i: t = f"<b><i>{t}</i></b>"
            elif b:     t = f"<b>{t}</b>"
            elif i:     t = f"<i>{t}</i>"
            if u:       t = f"<u>{t}</u>"
            parts.append(t)
        return "".join(parts)

    # Iterar body en orden real (párrafos Y tablas intercalados)
    body_elem  = doc.element.body
    para_list  = doc.paragraphs
    table_list = doc.tables
    para_idx   = 0
    table_idx  = 0

    for child in body_elem:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "p":
            if para_idx >= len(para_list):
                para_idx += 1
                continue
            para = para_list[para_idx]
            para_idx += 1

            if not para.text.strip():
                story.append(Spacer(1, 4))
                continue

            style = make_style(para)
            try:
                story.append(Paragraph(rich_text(para), style))
            except Exception:
                story.append(Paragraph(_safe(para.text), style))

        elif tag == "tbl":
            if table_idx >= len(table_list):
                table_idx += 1
                continue
            tbl = table_list[table_idx]
            table_idx += 1

            cell_st = ParagraphStyle("tc", parent=base["Normal"],
                                     fontName=_fn(), fontSize=9, leading=12)
            rows_data = []
            for row in tbl.rows:
                row_cells = []
                for cell in row.cells:
                    parts = []
                    for cp in cell.paragraphs:
                        rt = rich_text(cp)
                        if rt.strip():
                            parts.append(rt)
                    content = "<br/>".join(parts) if parts else _safe(cell.text)
                    try:
                        row_cells.append(Paragraph(content, cell_st))
                    except Exception:
                        row_cells.append(cell.text)
                rows_data.append(row_cells)

            if not rows_data:
                continue

            col_count = max(len(r) for r in rows_data)
            col_w     = (pw - lm - rm) / col_count

            header_bg = colors.HexColor("#4361ee")
            try:
                shd = tbl.rows[0].cells[0]._tc.find(f".//{wqn('w:shd')}")
                if shd is not None:
                    fill = shd.get(wqn("w:fill"))
                    if fill and fill != "auto" and len(fill) == 6:
                        c = _color(fill)
                        if c:
                            header_bg = c
            except Exception:
                pass

            ts = TableStyle([
                ("BACKGROUND",    (0,0), (-1,0),  header_bg),
                ("TEXTCOLOR",     (0,0), (-1,0),  colors.white),
                ("FONTNAME",      (0,0), (-1,0),  _fn(bold=True)),
                ("FONTSIZE",      (0,0), (-1,-1), 9),
                ("ROWBACKGROUNDS",(0,1), (-1,-1),
                 [colors.white, colors.HexColor("#f4f5fb")]),
                ("GRID",          (0,0), (-1,-1), 0.5, colors.HexColor("#d0d4e8")),
                ("LEFTPADDING",   (0,0), (-1,-1), 6),
                ("RIGHTPADDING",  (0,0), (-1,-1), 6),
                ("TOPPADDING",    (0,0), (-1,-1), 4),
                ("BOTTOMPADDING", (0,0), (-1,-1), 4),
                ("VALIGN",        (0,0), (-1,-1), "MIDDLE"),
            ])
            story.append(Spacer(1, 6))
            story.append(Table(rows_data, colWidths=[col_w]*col_count,
                               style=ts, repeatRows=1, hAlign="LEFT"))
            story.append(Spacer(1, 6))

    if not story:
        story.append(Paragraph("(Documento vacío)", base["Normal"]))

    pdf.build(story)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# XLSX → PDF
# ─────────────────────────────────────────────────────────────────────────────

def xlsx_to_pdf(file_bytes):
    if not (XLSX_OK and RL_OK):
        raise RuntimeError("openpyxl y reportlab son necesarios")

    _ensure_fonts()
    wb  = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    buf = io.BytesIO()

    page    = rl_landscape(A4)
    lm = rm = 1.2*cm
    tm = bm = 1.5*cm
    usable  = page[0] - lm - rm

    pdf   = SimpleDocTemplate(buf, pagesize=page,
                              leftMargin=lm, rightMargin=rm,
                              topMargin=tm, bottomMargin=bm)
    base  = getSampleStyleSheet()
    story = []

    sheet_st = ParagraphStyle("sheet", parent=base["Normal"],
                               fontName=_fn(bold=True), fontSize=13,
                               textColor=colors.HexColor("#14161f"),
                               spaceBefore=8, spaceAfter=4)
    cell_st  = ParagraphStyle("cell", parent=base["Normal"],
                               fontName=_fn(), fontSize=8, leading=10)
    hdr_st   = ParagraphStyle("hdr", parent=base["Normal"],
                               fontName=_fn(bold=True), fontSize=8,
                               leading=10, textColor=colors.white)

    def fmt_val(v):
        if v is None:
            return ""
        if isinstance(v, float):
            return str(int(v)) if v == int(v) else f"{v:,.4f}".rstrip("0").rstrip(".")
        if hasattr(v, "strftime"):
            try:
                s = v.strftime("%d/%m/%Y %H:%M")
                return s.rstrip(" 00:00").rstrip(" 0:00")
            except Exception:
                return str(v)
        return _safe(str(v))

    for s_idx, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]
        story.append(Paragraph(_safe(sheet_name), sheet_st))
        story.append(HRFlowable(width="100%", thickness=1.5,
                                color=colors.HexColor("#4361ee"), spaceAfter=4))

        rows = list(ws.iter_rows(values_only=True))
        while rows and all(c is None for c in rows[-1]):
            rows.pop()

        if not rows:
            story.append(Paragraph("(hoja vacía)", base["Normal"]))
            if s_idx < len(wb.sheetnames) - 1:
                story.append(PageBreak())
            continue

        max_cols = max(
            (max((i+1 for i, c in enumerate(row) if c is not None), default=0)
             for row in rows), default=0)
        max_cols = min(max_cols, 20)
        if max_cols == 0:
            story.append(Paragraph("(sin datos)", base["Normal"]))
            continue

        col_w = usable / max_cols
        data  = []
        for r_idx, row in enumerate(rows):
            row_cells = []
            is_hdr = r_idx == 0
            for c_idx in range(max_cols):
                val  = row[c_idx] if c_idx < len(row) else None
                text = fmt_val(val)
                row_cells.append(Paragraph(text, hdr_st if is_hdr else cell_st))
            data.append(row_cells)

        header_bg = colors.HexColor("#4361ee")
        try:
            fc = list(ws.iter_rows(min_row=1, max_row=1))[0][0].fill.fgColor
            if fc and fc.type == "rgb" and fc.rgb and fc.rgb != "00000000":
                c = _color(fc.rgb[2:])
                if c:
                    header_bg = c
        except Exception:
            pass

        ts = TableStyle([
            ("BACKGROUND",    (0,0), (-1,0),  header_bg),
            ("FONTSIZE",      (0,0), (-1,-1), 8),
            ("ROWBACKGROUNDS",(0,1), (-1,-1),
             [colors.white, colors.HexColor("#f4f5fb")]),
            ("GRID",          (0,0), (-1,-1), 0.3, colors.HexColor("#d0d4e8")),
            ("LEFTPADDING",   (0,0), (-1,-1), 4),
            ("RIGHTPADDING",  (0,0), (-1,-1), 4),
            ("TOPPADDING",    (0,0), (-1,-1), 3),
            ("BOTTOMPADDING", (0,0), (-1,-1), 3),
            ("VALIGN",        (0,0), (-1,-1), "MIDDLE"),
        ])
        story.append(Table(data, colWidths=[col_w]*max_cols,
                           style=ts, repeatRows=1, hAlign="LEFT"))
        story.append(Spacer(1, 8))

        if s_idx < len(wb.sheetnames) - 1:
            story.append(PageBreak())

    if not story:
        story.append(Paragraph("(Documento vacío)", base["Normal"]))

    pdf.build(story)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# PPTX → PDF
# ─────────────────────────────────────────────────────────────────────────────

def pptx_to_pdf(file_bytes):
    if not (PPTX_OK and RL_OK):
        raise RuntimeError("python-pptx y reportlab son necesarios")

    _ensure_fonts()
    prs = Presentation(io.BytesIO(file_bytes))
    buf = io.BytesIO()

    page    = rl_landscape(A4)
    lm = rm = 1.5*cm
    tm = bm = 1.2*cm

    pdf   = SimpleDocTemplate(buf, pagesize=page,
                              leftMargin=lm, rightMargin=rm,
                              topMargin=tm, bottomMargin=bm)
    base     = getSampleStyleSheet()
    story    = []
    n_slides = len(prs.slides)

    snum_st  = ParagraphStyle("snum",   parent=base["Normal"],
                               fontName=_fn(), fontSize=7.5, spaceAfter=1,
                               textColor=colors.HexColor("#9099bb"))
    title_st = ParagraphStyle("stitle", parent=base["Normal"],
                               fontName=_fn(bold=True), fontSize=22,
                               textColor=colors.HexColor("#14161f"),
                               spaceBefore=2, spaceAfter=8, leading=28)
    sub_st   = ParagraphStyle("ssub",   parent=base["Normal"],
                               fontName=_fn(bold=True), fontSize=14,
                               textColor=colors.HexColor("#4361ee"),
                               spaceAfter=6, leading=18)
    body_st  = ParagraphStyle("sbody",  parent=base["Normal"],
                               fontName=_fn(), fontSize=11,
                               leading=15, spaceAfter=2,
                               textColor=colors.HexColor("#2a2d3e"))
    bul_st   = ParagraphStyle("sbullet",parent=base["Normal"],
                               fontName=_fn(), fontSize=11,
                               leading=15, spaceAfter=2,
                               textColor=colors.HexColor("#2a2d3e"),
                               leftIndent=12, firstLineIndent=-10)

    for i, slide in enumerate(prs.slides):
        story.append(Paragraph(f"Diapositiva {i+1} / {n_slides}", snum_st))

        title_text    = ""
        content_items = []

        for shape in sorted(slide.shapes,
                            key=lambda s: (getattr(s, "top", 0) or 0,
                                           getattr(s, "left", 0) or 0)):
            if not shape.has_text_frame:
                continue
            is_title_ph = False
            try:
                ph = shape.placeholder_format
                if ph and ph.type in (PP_PLACEHOLDER.TITLE,
                                      PP_PLACEHOLDER.CENTER_TITLE):
                    is_title_ph = True
            except Exception:
                pass

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                level = para.level or 0
                size, bold = 11.0, False
                try:
                    for run in para.runs:
                        if run.font.size:
                            size = run.font.size.pt
                            break
                    bold = any(r.font.bold for r in para.runs
                               if r.font.bold is not None)
                except Exception:
                    pass

                if is_title_ph and not title_text:
                    title_text = text
                else:
                    content_items.append((level, text, size, bold))

        if title_text:
            story.append(Paragraph(_safe(title_text), title_st))
            story.append(HRFlowable(width="100%", thickness=2,
                                    color=colors.HexColor("#4361ee"),
                                    spaceAfter=6))

        for level, text, size, bold in content_items:
            te = _safe(text)
            if level == 0 and not title_text:
                story.append(Paragraph(te, title_st))
                story.append(HRFlowable(width="100%", thickness=2,
                                        color=colors.HexColor("#4361ee"),
                                        spaceAfter=6))
                title_text = text
            elif level == 0 and bold and size >= 13:
                story.append(Paragraph(te, sub_st))
            elif level == 0:
                story.append(Paragraph(te, body_st))
            else:
                story.append(Paragraph("  " * (level-1) + f"\u2022 {te}", bul_st))

        if i < n_slides - 1:
            story.append(Spacer(1, 8))
            story.append(HRFlowable(width="100%", thickness=0.5,
                                    color=colors.HexColor("#e0e3f0"),
                                    spaceAfter=8))
            story.append(PageBreak())

    if not story:
        story.append(Paragraph("(Presentación vacía)", base["Normal"]))

    pdf.build(story)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# HTTP handler
# ─────────────────────────────────────────────────────────────────────────────

CONVERTERS = {
    ".docx": docx_to_pdf,
    ".xlsx": xlsx_to_pdf,
    ".pptx": pptx_to_pdf,
}


class handler(BaseHTTPRequestHandler):

    def do_OPTIONS(self):
        self._cors()
        self.end_headers()

    def do_POST(self):
        try:
            file_bytes, filename, err = _parse_multipart(self)
            if err:
                self._error(400, err)
                return
            if not filename:
                self._error(400, "No se recibió nombre de archivo")
                return
            ext = os.path.splitext(filename)[1].lower()
            if ext not in SUPPORTED:
                self._error(400,
                    f"Formato no soportado: '{ext}'. Usa .docx, .xlsx o .pptx")
                return
            if not file_bytes:
                self._error(400, "El archivo está vacío")
                return

            pdf_bytes = CONVERTERS[ext](file_bytes)
            out_name  = os.path.splitext(filename)[0] + ".pdf"

            self.send_response(200)
            self.send_header("Content-Type", "application/pdf")
            self.send_header("Content-Disposition",
                             f'attachment; filename="{out_name}"')
            self.send_header("Content-Length", str(len(pdf_bytes)))
            self._cors()
            self.end_headers()
            self.wfile.write(pdf_bytes)

        except Exception as e:
            self._error(500, f"Error al convertir: {str(e).splitlines()[0]}")

    def _cors(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def _error(self, code, msg):
        body = json.dumps({"error": msg}).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self._cors()
        self.end_headers()
        self.wfile.write(body)

    def log_message(self, fmt, *args):
        pass
