"""
API endpoint: POST /api/pdf_to_ppt
Convierte páginas de PDF a diapositivas PowerPoint editables.

Entorno: Vercel serverless.
Librerías: pdfplumber + python-pptx (100% Python, sin binarios del sistema)

Fixes v3:
- Bullets nativos de PowerPoint (buChar XML) en lugar de "• " como texto plano
- Detección de caracteres bullet en el PDF (•, -, *, ▪, ►, etc.)
- Detección de líneas horizontales del PDF (LTLine/LTRect) → líneas reales en PPT
- Indentación del texto como señal de lista
- Orden de extracción top→bottom corregido
"""
import io
import os
import cgi
import json
import re
from http.server import BaseHTTPRequestHandler

try:
    import pdfplumber
    PDFPLUMBER_OK = True
except ImportError:
    PDFPLUMBER_OK = False

try:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer, LTChar, LTLine, LTRect
    PDFMINER_OK = True
except ImportError:
    PDFMINER_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
    PPTX_OK = True
except ImportError:
    PPTX_OK = False


def _parse_multipart(handler):
    ct = handler.headers.get("Content-Type", "")
    if "multipart/form-data" not in ct:
        return None, None, "Content-Type must be multipart/form-data"
    env = {"REQUEST_METHOD": "POST", "CONTENT_TYPE": ct,
           "CONTENT_LENGTH": handler.headers.get("Content-Length", "0")}
    body = handler.rfile.read(int(handler.headers.get("Content-Length", 0)))
    fs = cgi.FieldStorage(fp=io.BytesIO(body), environ=env, keep_blank_values=True)
    if "file" not in fs:
        return None, None, "No 'file' field in request"
    f = fs["file"]
    return f.file.read(), f.filename, None


# ── Colores ───────────────────────────────────────────────────────────────────

C_ACCENT = RGBColor(0x43, 0x61, 0xEE)
C_DARK   = RGBColor(0x14, 0x16, 0x1F)
C_MUTED  = RGBColor(0x70, 0x79, 0x99)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BODY   = RGBColor(0x2A, 0x2D, 0x3E)
C_SUB    = RGBColor(0x43, 0x61, 0xEE)
C_BORDER = RGBColor(0xD8, 0xDB, 0xEC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Regex: caracteres que indican bullet al inicio de una línea
RE_BULLET = re.compile(
    r"^[\u2022\u2023\u2043\u25cf\u25e6\u25aa\u25ab\u25fb\u25fc"
    r"\u2043\u204c\u204d\u2219\u00b7\u25b8\u25b9\u25ba\u25bb"
    r"\u2192\u2794\u27a2\u2714\u2713\*\-\+]\s+"
)


# ── Helpers de forma ──────────────────────────────────────────────────────────

def _rect(slide, x, y, w, h, fill=None, line_color=None, line_w=0.5):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _textbox(slide, x, y, w, h, text, size, bold=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or C_DARK
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb


def _set_bg(slide, rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def _add_native_bullet(para, bullet_color_hex="4361EE"):
    """Activa bullet nativo de PowerPoint en un párrafo via XML."""
    pPr = para._p.get_or_add_pPr()
    # Eliminar buNone si existe
    for child in list(pPr):
        if child.tag in (qn("a:buNone"), qn("a:buChar")):
            pPr.remove(child)
    # Añadir carácter bullet
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "\u2022")
    # Sangría izquierda
    pPr.set("marL", "342900")
    pPr.set("indent", "-342900")
    # Color del bullet
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    srgb  = etree.SubElement(buClr, qn("a:srgbClr"))
    srgb.set("val", bullet_color_hex)


# ── Extracción de contenido ───────────────────────────────────────────────────

def _extract_page_content_plumber(page):
    """
    Extrae bloques de texto + detecta líneas horizontales.
    Devuelve lista de dicts: {text, x0, y0, x1, y1, size, bold, is_bullet, is_hr}
    Coordenadas normalizadas 0-1 (0,0 = esquina superior izquierda).
    """
    pw = float(page.width)
    ph = float(page.height)
    blocks = []

    # ── Detectar líneas/rectángulos horizontales ──────────────────────────
    hr_ys = {}  # norm_y → True (deduplicar)
    for collection in (getattr(page, "lines", []), getattr(page, "rects", [])):
        for obj in collection:
            try:
                x0, y0 = obj["x0"], obj["top"]
                x1, y1 = obj["x1"], obj["bottom"]
                w = abs(x1 - x0)
                h = abs(y1 - y0)
                if w > pw * 0.25 and h < 4:
                    hr_ys[round(y0 / ph, 3)] = True
            except Exception:
                pass

    for norm_y in sorted(hr_ys):
        blocks.append({
            "text": "", "x0": 0, "y0": norm_y, "x1": 1, "y1": norm_y,
            "size": 0, "bold": False, "is_bullet": False, "is_hr": True,
        })

    # ── Extraer palabras con atributos ────────────────────────────────────
    try:
        words = page.extract_words(
            x_tolerance=3, y_tolerance=3,
            keep_blank_chars=False,
            use_text_flow=True,
            extra_attrs=["size", "fontname"],
        )
    except Exception:
        words = page.extract_words(x_tolerance=3, y_tolerance=3)

    if not words:
        blocks.sort(key=lambda b: b["y0"])
        return blocks

    # Agrupar en líneas (top → bottom)
    lines = []
    curr  = []
    last_y = None
    for word in sorted(words, key=lambda w: (w["top"], w["x0"])):
        if last_y is None or abs(word["top"] - last_y) < 6:
            curr.append(word)
            last_y = word["top"]
        else:
            lines.append(curr)
            curr   = [word]
            last_y = word["top"]
    if curr:
        lines.append(curr)

    for line_words in lines:
        if not line_words:
            continue
        text = " ".join(w["text"] for w in line_words)
        x0   = min(w["x0"]     for w in line_words) / pw
        y0   = min(w["top"]    for w in line_words) / ph
        x1   = max(w["x1"]     for w in line_words) / pw
        y1   = max(w["bottom"] for w in line_words) / ph
        size = max(w.get("size", 11) for w in line_words)
        bold = any("Bold" in str(w.get("fontname", "")) for w in line_words)

        # Detectar bullet por carácter inicial
        is_bullet = bool(RE_BULLET.match(text))
        if is_bullet:
            text = RE_BULLET.sub("", text).strip()
            if not text:
                continue

        # Detectar bullet por indentación (x0 > 8% de la página)
        if not is_bullet and x0 > 0.08:
            is_bullet = True

        blocks.append({
            "text": text, "x0": x0, "y0": y0, "x1": x1, "y1": y1,
            "size": size, "bold": bold,
            "is_bullet": is_bullet, "is_hr": False,
        })

    blocks.sort(key=lambda b: b["y0"])
    return blocks


def _extract_page_content_pdfminer(page_layout):
    """Fallback con pdfminer — sin detección de HRs."""
    pw = float(page_layout.width)
    ph = float(page_layout.height)
    blocks = []

    for element in page_layout:
        if isinstance(element, (LTLine, LTRect)):
            try:
                w = abs(element.x1 - element.x0)
                h = abs(element.y1 - element.y0)
                if w > pw * 0.25 and h < 4:
                    norm_y = round(1 - element.y1 / ph, 3)
                    blocks.append({
                        "text": "", "x0": 0, "y0": norm_y, "x1": 1, "y1": norm_y,
                        "size": 0, "bold": False, "is_bullet": False, "is_hr": True,
                    })
            except Exception:
                pass

        elif isinstance(element, LTTextContainer):
            text = element.get_text().strip()
            if not text:
                continue
            size, bold = 11.0, False
            for line in element:
                for char in line:
                    if isinstance(char, LTChar):
                        size = char.size
                        bold = "Bold" in (char.fontname or "")
                        break
                break

            x0 = element.x0 / pw
            y0 = 1 - element.y1 / ph
            x1 = element.x1 / pw
            y1 = 1 - element.y0 / ph

            is_bullet = bool(RE_BULLET.match(text))
            if is_bullet:
                text = RE_BULLET.sub("", text).strip()
                if not text:
                    continue
            if not is_bullet and x0 > 0.08:
                is_bullet = True

            blocks.append({
                "text": text, "x0": x0, "y0": y0, "x1": x1, "y1": y1,
                "size": size, "bold": bold,
                "is_bullet": is_bullet, "is_hr": False,
            })

    blocks.sort(key=lambda b: b["y0"])
    return blocks


def _classify(blocks):
    """
    Clasifica bloques: título / subtítulo / cuerpo.
    Los HRs y bullets van siempre al cuerpo.
    """
    if not blocks:
        return "", "", []

    text_blocks = [b for b in blocks
                   if not b["is_hr"] and not b["is_bullet"] and b["text"].strip()]
    if not text_blocks:
        return "", "", list(blocks)

    max_size = max(b["size"] for b in text_blocks)

    title = subtitle = ""
    body  = []

    for b in sorted(blocks, key=lambda b: b["y0"]):
        if b["is_hr"] or b["is_bullet"]:
            body.append(b)
            continue
        text = b["text"].strip()
        if not text:
            continue

        if not title and (b["size"] >= max_size * 0.85 or
                          (b["bold"] and b["size"] >= max_size * 0.65)):
            title = text
        elif not subtitle and b["y0"] < 0.4 and b["size"] >= max_size * 0.55:
            subtitle = text
        else:
            body.append(b)

    return title, subtitle, body


# ── Construcción de diapositiva ───────────────────────────────────────────────

def _add_pptx_table(slide, table_data, x, y, w, h):
    if not table_data:
        return
    rows = len(table_data)
    cols = max(len(r) for r in table_data)
    shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl   = shape.table
    for c in range(cols):
        tbl.columns[c].width = w // cols
    for r_idx, row in enumerate(table_data):
        for c_idx in range(cols):
            val  = row[c_idx] if c_idx < len(row) else ""
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val) if val else ""
            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size  = Pt(9)
                    run.font.bold  = (r_idx == 0)
                    run.font.color.rgb = C_WHITE if r_idx == 0 else C_BODY
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_ACCENT
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF5, 0xFB)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE


def build_slide(prs, slide_num, n_slides, title, subtitle, body_blocks,
                page_tables=None):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))

    # Barra superior de acento
    _rect(slide, 0, 0, SLIDE_W, Inches(0.07), fill=C_ACCENT)

    # Cabecera gris claro
    _rect(slide, 0, Inches(0.07), SLIDE_W, Inches(1.25),
          fill=RGBColor(0xF7, 0xF8, 0xFC))

    # Badge número de diapositiva
    bw, bh = Inches(0.55), Inches(0.28)
    badge  = _rect(slide,
                   SLIDE_W - bw - Inches(0.28), Inches(0.16),
                   bw, bh, fill=C_ACCENT)
    badge.text_frame.text = f"{slide_num}/{n_slides}"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in badge.text_frame.paragraphs[0].runs:
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = C_WHITE

    # Título
    if title:
        _textbox(slide, Inches(0.55), Inches(0.16),
                 SLIDE_W - Inches(1.5), Inches(0.8),
                 title, 22, bold=True, color=C_DARK, align=PP_ALIGN.LEFT)

    # Subtítulo
    if subtitle:
        _textbox(slide, Inches(0.55), Inches(0.88),
                 SLIDE_W - Inches(1.2), Inches(0.4),
                 subtitle, 13, color=C_SUB, align=PP_ALIGN.LEFT)

    # Línea divisoria cabecera/cuerpo
    _rect(slide, Inches(0.55), Inches(1.33),
          SLIDE_W - Inches(1.1), Inches(0.015), fill=C_BORDER)

    margin = Inches(0.55)
    body_w = SLIDE_W - margin * 2
    body_y = Inches(1.45)
    body_h = SLIDE_H - body_y - Inches(0.35)

    # Tablas
    if page_tables:
        table_area_h = body_h * 0.5 if body_blocks else body_h
        table_area_h = max(table_area_h, Inches(1.5))
        for t_idx, tbl_data in enumerate(page_tables[:2]):
            t_h = table_area_h // len(page_tables[:2])
            _add_pptx_table(slide, tbl_data,
                            margin, body_y + t_idx * t_h,
                            body_w, t_h - Inches(0.1))
        body_y += table_area_h + Inches(0.1)
        body_h  = SLIDE_H - body_y - Inches(0.35)

    # ── Cuerpo: texto, bullets nativos y separadores reales ───────────────
    if body_blocks and body_h > Inches(0.25):
        # Partir en grupos separados por HRs
        groups  = []
        current = []
        for b in body_blocks:
            if b["is_hr"]:
                if current:
                    groups.append(("text", current))
                    current = []
                groups.append(("hr", None))
            else:
                current.append(b)
        if current:
            groups.append(("text", current))

        n_text  = sum(1 for k, _ in groups if k == "text")
        n_hr    = sum(1 for k, _ in groups if k == "hr")
        hr_h    = Inches(0.03)
        hr_gap  = Inches(0.1)
        avail_h = body_h - n_hr * (hr_h + hr_gap * 2)
        text_h  = avail_h / max(n_text, 1)

        cur_y = body_y

        for kind, data in groups:
            if kind == "hr":
                _rect(slide,
                      margin + Inches(0.15),
                      cur_y + hr_gap,
                      body_w - Inches(0.3),
                      hr_h,
                      fill=C_BORDER)
                cur_y += hr_h + hr_gap * 2

            else:
                if not data:
                    continue
                tb = slide.shapes.add_textbox(margin, cur_y, body_w, text_h)
                tf = tb.text_frame
                tf.word_wrap = True

                first = True
                for b in data:
                    text = b["text"].strip()
                    if not text:
                        continue

                    size      = b["size"]
                    bold      = b["bold"]
                    is_bullet = b["is_bullet"]
                    is_heading = bold and size >= 13 and not is_bullet

                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.alignment = PP_ALIGN.LEFT

                    if is_bullet:
                        _add_native_bullet(p)

                    run = p.add_run()
                    run.text = text
                    run.font.size  = Pt(min(max(float(size), 9), 20))
                    run.font.bold  = is_heading
                    run.font.color.rgb = C_DARK if is_heading else C_BODY
                    p.space_after = Pt(4 if is_heading else 2)

                cur_y += text_h

    # Pie de página
    _rect(slide, 0, SLIDE_H - Inches(0.28), SLIDE_W, Inches(0.015),
          fill=C_BORDER)

    return slide


# ── Conversión principal ──────────────────────────────────────────────────────

def _extract_tables_plumber(page):
    for settings in [
        {"vertical_strategy": "lines_strict",
         "horizontal_strategy": "lines_strict", "snap_tolerance": 4},
        {"vertical_strategy": "lines",
         "horizontal_strategy": "lines", "snap_tolerance": 5},
    ]:
        try:
            tables = page.extract_tables(settings)
            valid  = [t for t in (tables or [])
                      if any(any(c for c in row if c) for row in t)]
            if valid:
                return valid
        except Exception:
            pass
    return []


def pdf_to_pptx(file_bytes):
    if not PPTX_OK:
        raise RuntimeError("python-pptx no está instalado")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    if PDFPLUMBER_OK:
        pages_data = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                blocks = _extract_page_content_plumber(page)
                tables = _extract_tables_plumber(page)
                pages_data.append((blocks, tables))
    elif PDFMINER_OK:
        pages_data = []
        for page_layout in extract_pages(io.BytesIO(file_bytes)):
            blocks = _extract_page_content_pdfminer(page_layout)
            pages_data.append((blocks, []))
    else:
        raise RuntimeError("Instala pdfplumber: pip install pdfplumber python-pptx")

    n_slides = len(pages_data)

    for idx, (blocks, tables) in enumerate(pages_data):
        title, subtitle, body = _classify(blocks)
        build_slide(prs, idx + 1, n_slides, title, subtitle, body, tables)

    if not pages_data:
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
        _textbox(slide, SLIDE_W // 4, SLIDE_H // 3,
                 SLIDE_W // 2, Inches(1),
                 "PDF sin contenido extraíble", 18,
                 bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── HTTP handler ──────────────────────────────────────────────────────────────

class handler(BaseHTTPRequestHandler):

    def do_OPTIONS(self):
        self._cors()
        self.end_headers()

    def do_POST(self):
        try:
            file_bytes, filename, err = _parse_multipart(self)
            if err:
                self._error(400, err)
                return
            if not filename or not filename.lower().endswith(".pdf"):
                self._error(400, "Por favor sube un archivo PDF (.pdf)")
                return
            if not file_bytes:
                self._error(400, "El archivo está vacío")
                return

            pptx_bytes = pdf_to_pptx(file_bytes)
            out_name   = filename.rsplit(".", 1)[0] + ".pptx"

            self.send_response(200)
            self.send_header(
                "Content-Type",
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation",
            )
            self.send_header("Content-Disposition",
                             f'attachment; filename="{out_name}"')
            self.send_header("Content-Length", str(len(pptx_bytes)))
            self._cors()
            self.end_headers()
            self.wfile.write(pptx_bytes)

        except Exception as e:
            self._error(500, f"Error al convertir: {str(e).splitlines()[0]}")

    def _cors(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def _error(self, code, msg):
        body = json.dumps({"error": msg}).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self._cors()
        self.end_headers()
        self.wfile.write(body)

    def log_message(self, fmt, *args):
        pass
